#!/usr/bin/env python3
"""
Create a comprehensive PowerPoint presentation for the Tissue QA Dashboard
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def add_title_slide(prs):
    """Add title slide"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Tissue Manufacturing QA Dashboard"
    subtitle.text = "Comprehensive Quality Analytics & Monitoring System\nReal-time Process Control & Performance Insights"
    
    # Style the title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    return slide

def add_overview_slide(prs):
    """Add overview slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Dashboard Overview"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Complete Quality Management Solution"
    
    points = [
        "Real-time quality monitoring across all production parameters",
        "Advanced analytics with daily averaging for stable insights",
        "Multi-dimensional filtering and data exploration",
        "Automated PDF report generation",
        "Process capability analysis (Cpk)",
        "Trend analysis with statistical control limits",
        "Shift performance comparison",
        "Correlation analysis between parameters"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(18)
    
    return slide

def add_key_features_slide(prs):
    """Add key features slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Key Features"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. Data Upload & Processing"
    tf.paragraphs[0].font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Excel file upload with automatic parsing"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Intelligent column mapping"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Data validation and error handling"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n2. Daily Quality Reports"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Comprehensive daily metrics view"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Automatic data aggregation"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• PDF report generation"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n3. Advanced Filtering"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Date range selection"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Shift, Quality, and GSM grade filters"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Multi-criteria filtering"
    p.level = 1
    
    return slide

def add_metrics_monitored_slide(prs):
    """Add metrics monitored slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Quality Metrics Monitored"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Core Quality Parameters"
    tf.paragraphs[0].font.bold = True
    
    metrics = [
        "GSM (Grammage) - g/m²",
        "Thickness - μm",
        "Bulk - cc/g",
        "Tensile Strength MD/CD - N/m",
        "MD/CD Ratio",
        "Brightness ISO - %",
        "Opacity - %",
        "Moisture Content - %",
        "Stretch/Elongation - %",
        "Wet Tensile - gf/50mm",
        "Wet/Dry Tensile Ratio - %"
    ]
    
    for metric in metrics:
        p = tf.add_paragraph()
        p.text = f"• {metric}"
        p.level = 1
        p.font.size = Pt(16)
    
    return slide

def add_machine_parameters_slide(prs):
    """Add machine parameters slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Machine Parameters & Fiber Composition"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Machine Parameters"
    tf.paragraphs[0].font.bold = True
    
    params = [
        "Machine Speed (Mpm)",
        "Pope Reel Speed (Mpm)",
        "MC Draw",
        "Press Load",
        "Coating Parameters",
        "Machine Creep %"
    ]
    
    for param in params:
        p = tf.add_paragraph()
        p.text = f"• {param}"
        p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nFiber Composition & Consumption"
    p.font.bold = True
    
    fiber_params = [
        "Short Fiber %",
        "Long Fiber %",
        "Broke %",
        "HW/SW Consistency",
        "HW SR / SW OSR",
        "WSR/DSR (Kg/Hr)"
    ]
    
    for param in fiber_params:
        p = tf.add_paragraph()
        p.text = f"• {param}"
        p.level = 1
    
    return slide

def add_trend_analysis_slide(prs):
    """Add trend analysis features slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Trend Analysis Features"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Advanced Trending Capabilities"
    tf.paragraphs[0].font.bold = True
    
    features = [
        "Multiple time views: Hourly, Daily, Weekly, Monthly",
        "Multi-metric selection (up to 4 parameters)",
        "Statistical indicators:",
        "  - Moving averages (customizable period)",
        "  - Control limits (3-sigma)",
        "  - Min/Max/Mean/Median statistics",
        "Interactive charts with zoom and pan",
        "Chart export functionality",
        "CSV data export",
        "Filter selections shown in exports"
    ]
    
    for feature in features:
        p = tf.add_paragraph()
        p.text = feature
        p.level = 1 if not feature.startswith("  ") else 2
        p.font.size = Pt(16)
    
    return slide

def add_advanced_analytics_slide(prs):
    """Add advanced analytics slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Advanced Analytics Dashboard"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Four Comprehensive Analysis Tabs"
    tf.paragraphs[0].font.bold = True
    
    tabs = [
        "\n1. Process Performance",
        "  • Process Capability Index (Cpk) analysis",
        "  • Multi-parameter performance radar",
        "  • Quality score trending",
        "  • Daily averaging for stability",
        
        "\n2. Statistical Analysis",
        "  • Process stability monitoring",
        "  • Distribution analysis",
        "  • Coefficient of variation",
        "  • Control charts",
        
        "\n3. Correlations & Patterns",
        "  • Parameter correlation matrix",
        "  • Top quality issues distribution",
        "  • Pattern recognition",
        
        "\n4. Shift Performance",
        "  • Comparative shift analysis",
        "  • Performance benchmarking",
        "  • Production volume tracking"
    ]
    
    for item in tabs:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0 if item.startswith("\n") else 1
        p.font.size = Pt(16)
        if item.startswith("\n"):
            p.font.bold = True
    
    return slide

def add_insights_slide1(prs):
    """Add insights slide 1"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Key Insights & Benefits"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Quality Control Excellence"
    tf.paragraphs[0].font.bold = True
    
    insights = [
        "Real-time identification of out-of-spec parameters",
        "Early warning system for process deviations",
        "Cpk > 1.33 indicates excellent process capability",
        "Daily averaging reduces noise in measurements",
        
        "\nOperational Efficiency",
        "Shift performance comparison identifies best practices",
        "Machine parameter tracking optimizes settings",
        "Correlation analysis reveals parameter relationships",
        
        "\nData-Driven Decision Making",
        "Historical trend analysis for predictive insights",
        "Statistical control limits for process stability",
        "Comprehensive reporting for management review"
    ]
    
    for insight in insights:
        p = tf.add_paragraph()
        p.text = insight
        p.level = 0 if insight.startswith("\n") else 1
        p.font.size = Pt(16)
        if insight.startswith("\n"):
            p.font.bold = True
    
    return slide

def add_insights_slide2(prs):
    """Add insights slide 2"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Process Optimization Insights"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Critical Parameters to Monitor"
    tf.paragraphs[0].font.bold = True
    
    insights = [
        "Moisture content (exclude 0 values for accuracy)",
        "MD/CD ratio for sheet strength balance",
        "Wet/Dry tensile ratio for absorbency",
        "Process stability through control charts",
        
        "\nQuality Improvement Opportunities",
        "Parameters frequently out of spec",
        "Shift-to-shift variations",
        "Correlation between defects and parameters",
        "Machine speed vs quality trade-offs",
        
        "\nCost Optimization",
        "Fiber composition optimization",
        "Energy consumption (WSR/DSR rates)",
        "Broke percentage reduction",
        "Machine efficiency improvements"
    ]
    
    for insight in insights:
        p = tf.add_paragraph()
        p.text = insight
        p.level = 0 if insight.startswith("\n") else 1
        p.font.size = Pt(16)
        if insight.startswith("\n"):
            p.font.bold = True
    
    return slide

def add_technical_features_slide(prs):
    """Add technical features slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Technical Implementation"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Modern Technology Stack"
    tf.paragraphs[0].font.bold = True
    
    features = [
        "React 19 with TypeScript",
        "Material-UI v7 for modern UI",
        "Recharts for interactive visualizations",
        "Day.js for date handling",
        "XLSX for Excel file parsing",
        "jsPDF for report generation",
        
        "\nData Processing Features",
        "Automatic column mapping",
        "Intelligent date parsing",
        "Data validation and error handling",
        "Daily averaging algorithms",
        "Statistical calculations",
        
        "\nUser Experience",
        "Responsive design for all devices",
        "Intuitive navigation",
        "Real-time filtering",
        "Export capabilities",
        "Performance optimized"
    ]
    
    for feature in features:
        p = tf.add_paragraph()
        p.text = feature
        p.level = 0 if feature.startswith("\n") else 1
        p.font.size = Pt(16)
        if feature.startswith("\n"):
            p.font.bold = True
    
    return slide

def add_future_enhancements_slide(prs):
    """Add future enhancements slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Future Enhancement Opportunities"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Potential Additions"
    tf.paragraphs[0].font.bold = True
    
    enhancements = [
        "Machine Learning Integration",
        "  • Predictive quality alerts",
        "  • Anomaly detection",
        "  • Optimal parameter recommendations",
        
        "\nReal-time Data Integration",
        "  • Direct sensor connectivity",
        "  • Live dashboard updates",
        "  • Instant notifications",
        
        "\nAdvanced Analytics",
        "  • Six Sigma calculations",
        "  • Root cause analysis",
        "  • Predictive maintenance",
        
        "\nIntegration Capabilities",
        "  • ERP system integration",
        "  • Mobile app development",
        "  • API for third-party access"
    ]
    
    for item in enhancements:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0 if not item.startswith("  ") else 1
        p.font.size = Pt(16)
        if not item.startswith("  ") and not item == "Potential Additions":
            p.font.bold = True
    
    return slide

def add_conclusion_slide(prs):
    """Add conclusion slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Conclusion"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Comprehensive Quality Management Solution"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)
    
    points = [
        "\nKey Benefits:",
        "• Improved quality control and consistency",
        "• Data-driven decision making",
        "• Reduced waste and defects",
        "• Enhanced operational efficiency",
        "• Better compliance and reporting",
        
        "\nBusiness Impact:",
        "• Faster identification of quality issues",
        "• Improved customer satisfaction",
        "• Cost reduction through optimization",
        "• Competitive advantage through analytics",
        
        "\nReady for deployment and continuous improvement"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0 if point.startswith("\n") else 1
        p.font.size = Pt(18)
        if point.startswith("\n"):
            p.font.bold = True
    
    return slide

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    
    # Set slide size to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Add all slides
    add_title_slide(prs)
    add_overview_slide(prs)
    add_key_features_slide(prs)
    add_metrics_monitored_slide(prs)
    add_machine_parameters_slide(prs)
    add_trend_analysis_slide(prs)
    add_advanced_analytics_slide(prs)
    add_insights_slide1(prs)
    add_insights_slide2(prs)
    add_technical_features_slide(prs)
    add_future_enhancements_slide(prs)
    add_conclusion_slide(prs)
    
    # Save presentation
    filename = "Tissue_QA_Dashboard_Presentation.pptx"
    prs.save(filename)
    print(f"Presentation saved as {filename}")
    
    return filename

if __name__ == "__main__":
    create_presentation()